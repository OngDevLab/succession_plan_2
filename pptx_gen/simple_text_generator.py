"""
Simple PowerPoint generator - add plain text after headers, left-aligned, no expansion
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
import io
from urllib.request import urlopen
from PIL import Image, ImageDraw
from config.loader import CONFIG
from utils.pptx_repair import auto_repair_pptx

def create_succession_plan_from_template(incumbent_data, successors_data):
    """Create PowerPoint with multiple slides if needed (configurable successors per slide)"""
    
    # Load the template
    prs = Presentation(CONFIG['powerpoint']['template_file'])
    
    # Split successors into groups (max per slide from config)
    successors_per_slide = CONFIG['powerpoint']['successors_per_slide']
    successor_groups = []
    for i in range(0, len(successors_data), successors_per_slide):
        successor_groups.append(successors_data[i:i + successors_per_slide])
    
    print(f"Creating {len(successor_groups)} slides for {len(successors_data)} successors")
    
    # For additional slides, we need to duplicate the template slide completely
    original_slide = prs.slides[0]
    
    # Add additional slides by duplicating the entire template slide
    for group_idx in range(1, len(successor_groups)):
        # Create a new presentation with just the template to copy from
        temp_prs = Presentation(CONFIG['powerpoint']['template_file'])
        template_slide = temp_prs.slides[0]
        
        # Get the slide layout
        slide_layout = template_slide.slide_layout
        
        # Add new slide with same layout
        new_slide = prs.slides.add_slide(slide_layout)
        
        # Copy all elements from template slide to new slide
        copy_slide_elements(template_slide, new_slide)
    
    # Process each group on its respective slide
    for group_idx, successor_group in enumerate(successor_groups):
        slide = prs.slides[group_idx]
        print(f"Processing slide {group_idx + 1} with {len(successor_group)} successors")
        fill_template_simple_text(slide, incumbent_data, successor_group)
    
    # Save to BytesIO
    pptx_buffer = io.BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    
    # Auto-repair the PowerPoint file if enabled
    if CONFIG['powerpoint'].get('auto_repair', True):
        repair_method = CONFIG['powerpoint'].get('repair_method', 'standard')
        print(f"🔧 Auto-repairing PowerPoint file using {repair_method} method...")
        repaired_buffer = auto_repair_pptx(pptx_buffer, method=repair_method)
        return repaired_buffer
    else:
        print("⚠️ Auto-repair disabled - returning original file")
        return pptx_buffer

def copy_slide_elements(source_slide, target_slide):
    """Copy all elements from source slide to target slide"""
    try:
        # Import the XML copying functionality
        from lxml import etree
        import copy
        
        # Get the source slide's shape tree
        source_tree = source_slide._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}spTree')
        target_tree = target_slide._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}spTree')
        
        # Copy all shape elements except the first two (which are slide properties)
        for shape_elem in list(source_tree)[2:]:  # Skip first 2 elements (nvGrpSpPr and grpSpPr)
            # Create a deep copy of the shape element
            new_shape_elem = copy.deepcopy(shape_elem)
            target_tree.append(new_shape_elem)
            
        print(f"Copied {len(list(source_tree)[2:])} elements to new slide")
        
    except Exception as e:
        print(f"Error copying slide elements: {e}")
        # Fallback: try to copy basic shapes
        try:
            for shape in source_slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    # Copy text boxes
                    new_shape = target_slide.shapes.add_textbox(
                        shape.left, shape.top, shape.width, shape.height
                    )
                    new_shape.text = shape.text
                    # Copy basic formatting
                    for i, paragraph in enumerate(shape.text_frame.paragraphs):
                        if i < len(new_shape.text_frame.paragraphs):
                            new_shape.text_frame.paragraphs[i].alignment = paragraph.alignment
        except Exception as e2:
            print(f"Fallback copying also failed: {e2}")

def fill_template_simple_text(slide, incumbent_data, successors_data):
    """Fill template by replacing carrot placeholders only, FORCE LEFT ALIGNMENT"""
    
    incumbent_name = f"{incumbent_data['metadata']['PREFERRED_NAME_FIRST_NAME']} {incumbent_data['metadata']['PREFERRED_NAME_LAST_NAME']}"
    plan = incumbent_data['plan_details']
    
    shapes = list(slide.shapes)
    
    # Replace only specific placeholders, FORCE LEFT ALIGNMENT
    for shape in shapes:
        if hasattr(shape, 'text_frame') and shape.text_frame:
            # FORCE LEFT ALIGNMENT on all paragraphs
            for paragraph in shape.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT  # FORCE LEFT ALIGNMENT
                
                for run in paragraph.runs:
                    # Only replace specific placeholders
                    if "POSITION" == run.text.strip():
                        run.text = incumbent_data['metadata'].get('POSITION_NBR_DESCRIPTION', 'POSITION')
                    elif "NAME" == run.text.strip():
                        run.text = incumbent_name
                    elif "Insert role information summary" in run.text:
                        # Replace with proper incumbent summary like app_final.py
                        incumbent_summary = get_incumbent_summary_like_app_final(incumbent_data, plan)
                        run.text = incumbent_summary
                    elif run.text.strip() == "<you can edit this too> Focus on \"make or break\" descriptors (responsibilities, typical challenges, unique capabilities, qualities, and track record required for success in the role)":
                        # Replace with actual responsibilities
                        responsibilities = plan.get('responsibilities', 'Role responsibilities not specified')
                        run.text = responsibilities
                    elif run.text.strip() == "<This area you can add text>":
                        # Replace individual carrot placeholders with incumbent details
                        incumbent_details = get_incumbent_details_list(plan)
                        if incumbent_details:
                            run.text = incumbent_details.pop(0)
                        else:
                            run.text = ""  # Remove placeholder if no content
    
    # Fill the main table, FORCE LEFT ALIGNMENT
    if len(shapes) > 2 and shapes[2].shape_type == MSO_SHAPE_TYPE.TABLE:
        fill_table_carrot_placeholders_only(shapes[2].table, successors_data)
    
    # Replace photos with circular faces
    photo_shapes = shapes[5:9] if len(shapes) > 8 else []
    replace_with_circular_faces(slide, photo_shapes, incumbent_data, successors_data)

def get_incumbent_summary_like_app_final(incumbent_data, plan):
    """Get incumbent summary exactly like app_final.py display"""
    
    person = incumbent_data['metadata']
    full_name = f"{person['PREFERRED_NAME_FIRST_NAME']} {person['PREFERRED_NAME_LAST_NAME']}"
    
    # Build summary like app_final.py lines 552-565
    summary_parts = []
    
    # Name and position
    summary_parts.append(f"INCUMBENT: {full_name}")
    summary_parts.append(f"Position: {person['POSITION_NBR_DESCRIPTION']}")
    
    # Critical role
    critical_text = "Yes" if plan.get('critical_role') else "No"
    summary_parts.append(f"Critical Role: {critical_text}")
    
    # Sourcing strategy
    sourcing = plan.get('sourcing_strategy', 'N/A')
    summary_parts.append(f"Sourcing Strategy: {sourcing}")
    
    # Scenario
    scenario = plan.get('scenario_plan', 'N/A')
    summary_parts.append(f"Scenario: {scenario}")
    
    # Top PLE
    top_ple = plan.get('top_ple', 'Not specified')
    if top_ple != "-- Select an Option --":
        summary_parts.append(f"Top Demonstrated PLE: {top_ple}")
    
    # Top Skills
    top_skills = plan.get('top_skills', [])
    if top_skills:
        skills_text = ", ".join(top_skills)
        summary_parts.append(f"Top Skills: {skills_text}")
    
    # Responsibilities
    responsibilities = plan.get('responsibilities', 'Not specified')
    if responsibilities and responsibilities != 'Not specified':
        summary_parts.append(f"Responsibilities & Attributes: {responsibilities}")
    
    return "\n".join(summary_parts)

def get_incumbent_details_list(plan):
    """Get incumbent details exactly like app_final.py display"""
    details = []
    
    # Use exact same format as app_final.py lines 555-556
    critical_text = "Yes" if plan.get('critical_role') else "No"
    details.append(f"Critical Role: {critical_text}")
    details.append(f"Sourcing Strategy: {plan.get('sourcing_strategy', 'N/A')}")
    details.append(f"Scenario: {plan.get('scenario_plan', 'N/A')}")
    
    # Add top skills and PLE like app_final.py
    if plan.get('top_skills'):
        skills_text = ", ".join(plan.get('top_skills', []))
        details.append(f"Top Skills: {skills_text}")
    
    if plan.get('top_ple') and plan.get('top_ple') != "-- Select an Option --":
        details.append(f"Top PLE: {plan.get('top_ple')}")
    
    return details

def fill_table_carrot_placeholders_only(table, successors_data):
    """Fill table by replacing ONLY carrot placeholders, DON'T TOUCH HEADERS"""
    
    # Fill each column with successor data
    max_successors = CONFIG['powerpoint']['successors_per_slide']
    for col_idx, successor in enumerate(successors_data[:max_successors]):  # Max successors per slide
        if col_idx >= len(table.columns):
            break
            
        succ_name = f"{successor['metadata']['PREFERRED_NAME_FIRST_NAME']} {successor['metadata']['PREFERRED_NAME_LAST_NAME']}"
        succ_title = successor['metadata'].get('POSITION_NBR_DESCRIPTION', '')
        assessment = successor['assessment']
        
        # Row 0: Replace template placeholders with name/title/readiness
        if len(table.rows) > 0 and col_idx < len(table.rows[0].cells):
            cell = table.rows[0].cells[col_idx]
            readiness_text = f"{succ_name}\n{succ_title}\nReadiness: {assessment.get('readiness', '')}"
            if assessment.get('future_readiness_timing'):
                readiness_text += f" ({assessment.get('future_readiness_timing')})"
            
            # Replace all content (this cell doesn't have carrot placeholders)
            cell.text = readiness_text
            # FORCE LEFT ALIGNMENT on successor names/positions/readiness
            if cell.text_frame:
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.LEFT
        
        # Row 1: Replace ONLY carrot placeholders in strengths section - DON'T TOUCH HEADERS
        if len(table.rows) > 1 and col_idx < len(table.rows[1].cells):
            cell = table.rows[1].cells[col_idx]
            replace_carrot_placeholders_simple(cell, assessment, 'strengths')
        
        # Row 2: Replace ONLY carrot placeholders in development section - DON'T TOUCH HEADERS
        if len(table.rows) > 2 and col_idx < len(table.rows[2].cells):
            cell = table.rows[2].cells[col_idx]
            replace_carrot_placeholders_simple(cell, assessment, 'development')
        
        # Row 3: Replace ONLY carrot placeholders in talent actions section - DON'T TOUCH HEADERS
        if len(table.rows) > 3 and col_idx < len(table.rows[3].cells):
            cell = table.rows[3].cells[col_idx]
            replace_carrot_placeholders_simple(cell, assessment, 'actions')
    
    # Clear unused columns - FORCE LEFT ALIGNMENT on Row 0
    max_successors = CONFIG['powerpoint']['successors_per_slide']
    for col_idx in range(len(successors_data), max_successors):
        if col_idx < len(table.columns):
            for row_idx in range(len(table.rows)):
                if col_idx < len(table.rows[row_idx].cells):
                    cell = table.rows[row_idx].cells[col_idx]
                    if row_idx == 0:
                        cell.text = ""  # Clear name/title/readiness completely
                        # FORCE LEFT ALIGNMENT on Row 0
                        if cell.text_frame:
                            for paragraph in cell.text_frame.paragraphs:
                                paragraph.alignment = PP_ALIGN.LEFT
                    else:
                        # Just clear carrot placeholders, DON'T TOUCH HEADER ALIGNMENT
                        clear_carrot_placeholders_keep_headers(cell)

def replace_carrot_placeholders_simple(cell, assessment, content_type):
    """Replace carrot placeholders simply - FORCE LEFT ALIGNMENT"""
    
    if not cell.text_frame:
        return
    
    # Get content based on type
    if content_type == 'strengths':
        content_text = assessment.get('strengths', '')
        # Add skills if we have them
        top_skills = assessment.get('top_skills', [])
        if top_skills:
            skills_text = f"Skills: {', '.join(top_skills[:3])}"
            content_text = f"{content_text}\n{skills_text}" if content_text else skills_text
    elif content_type == 'development':
        content_text = assessment.get('development_focus', '')
    elif content_type == 'actions':
        content_text = assessment.get('talent_actions', '')
    else:
        content_text = ""
    
    # Split content into lines if needed
    if content_text:
        # Simple line splitting - don't overcomplicate
        lines = content_text.split('\n')
        # If single long line, break it up
        if len(lines) == 1 and len(content_text) > 80:
            words = content_text.split()
            lines = []
            current_line = ""
            for word in words:
                if len(current_line + " " + word) <= 80:
                    current_line += (" " + word) if current_line else word
                else:
                    if current_line:
                        lines.append(current_line)
                    current_line = word
            if current_line:
                lines.append(current_line)
    else:
        lines = []
    
    # Replace carrot placeholders with content lines - FORCE LEFT ALIGNMENT
    line_index = 0
    for paragraph in cell.text_frame.paragraphs:
        # FORCE LEFT ALIGNMENT
        paragraph.alignment = PP_ALIGN.LEFT
        
        for run in paragraph.runs:
            if run.text.strip() == "<This area you can add text>":
                if line_index < len(lines):
                    run.text = lines[line_index]
                    line_index += 1
                else:
                    run.text = ""  # Remove extra placeholders

def clear_carrot_placeholders_keep_headers(cell):
    """Clear carrot placeholders but keep headers and FORCE LEFT ALIGNMENT"""
    
    if not cell.text_frame:
        return
    
    for paragraph in cell.text_frame.paragraphs:
        # FORCE LEFT ALIGNMENT
        paragraph.alignment = PP_ALIGN.LEFT
        
        for run in paragraph.runs:
            if run.text.strip() == "<This area you can add text>":
                run.text = ""  # Remove placeholder

def replace_with_circular_faces(slide, photo_shapes, incumbent_data, successors_data):
    """Replace photo placeholders with circular cropped faces using config URL"""
    
    max_successors = CONFIG['powerpoint']['successors_per_slide']
    employees = [incumbent_data] + successors_data[:max_successors]
    
    for i, shape in enumerate(photo_shapes[:len(employees)]):
        if i < len(employees):
            employee = employees[i]
            employee_id = employee['metadata']['EMPLOYEE_ID']
            
            try:
                # Use URL from config
                photo_url = CONFIG['avatar']['url_template'].format(employee_id=employee_id)
                with urlopen(photo_url) as response:
                    photo_data = response.read()
                
                # Create circular image using PIL
                circular_photo_data = create_circular_image(photo_data)
                
                # Get original position and size
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                
                # Remove placeholder
                shape_element = shape.element
                shape_element.getparent().remove(shape_element)
                
                # Add circular photo
                circular_photo_stream = io.BytesIO(circular_photo_data)
                new_pic = slide.shapes.add_picture(circular_photo_stream, left, top, width, height)
                
            except Exception as e:
                print(f"Could not create circular photo for employee {employee_id}: {e}")
                # Keep placeholder

def create_circular_image(image_data):
    """Create a circular cropped image using PIL"""
    try:
        # Open image with PIL
        image = Image.open(io.BytesIO(image_data))
        
        # Convert to RGB if necessary
        if image.mode != 'RGB':
            image = image.convert('RGB')
        
        # Make it square by cropping to center
        width, height = image.size
        size = min(width, height)
        
        # Calculate crop box to center the image
        left = (width - size) // 2
        top = (height - size) // 2
        right = left + size
        bottom = top + size
        
        # Crop to square
        image = image.crop((left, top, right, bottom))
        
        # Create circular mask
        mask = Image.new('L', (size, size), 0)
        draw = ImageDraw.Draw(mask)
        draw.ellipse((0, 0, size, size), fill=255)
        
        # Apply circular mask
        output = Image.new('RGBA', (size, size), (0, 0, 0, 0))
        output.paste(image, (0, 0))
        output.putalpha(mask)
        
        # Save to bytes
        output_buffer = io.BytesIO()
        output.save(output_buffer, format='PNG')
        output_buffer.seek(0)
        
        return output_buffer.getvalue()
        
    except Exception as e:
        print(f"Error creating circular image: {e}")
        # Return original image data if circular creation fails
        return image_data
