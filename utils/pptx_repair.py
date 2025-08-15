"""
PowerPoint repair utilities - Clean up corrupted PPTX files
"""

import io
from pptx import Presentation
import tempfile
import os

def repair_pptx_buffer(pptx_buffer):
    """
    Repair a PowerPoint file by loading and re-saving it
    This mimics what PowerPoint's repair function does
    """
    try:
        print("🔧 Starting PowerPoint repair process...")
        
        # Reset buffer position
        pptx_buffer.seek(0)
        
        # Load the potentially corrupted presentation
        prs = Presentation(pptx_buffer)
        print(f"✅ Successfully loaded presentation with {len(prs.slides)} slides")
        
        # Create a new clean buffer
        repaired_buffer = io.BytesIO()
        
        # Save to new buffer (this cleans up the internal structure)
        prs.save(repaired_buffer)
        repaired_buffer.seek(0)
        
        print("✅ PowerPoint repair completed successfully")
        return repaired_buffer
        
    except Exception as e:
        print(f"❌ PowerPoint repair failed: {e}")
        # Return original buffer if repair fails
        pptx_buffer.seek(0)
        return pptx_buffer

def repair_pptx_with_temp_file(pptx_buffer):
    """
    Alternative repair method using temporary files
    Sometimes more effective for severely corrupted files
    """
    try:
        print("🔧 Starting advanced PowerPoint repair with temp files...")
        
        # Reset buffer position
        pptx_buffer.seek(0)
        
        # Create temporary files
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_input:
            temp_input.write(pptx_buffer.read())
            temp_input_path = temp_input.name
        
        try:
            # Load from temp file
            prs = Presentation(temp_input_path)
            print(f"✅ Loaded presentation with {len(prs.slides)} slides")
            
            # Save to another temp file
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_output:
                temp_output_path = temp_output.name
            
            prs.save(temp_output_path)
            
            # Read the repaired file back into buffer
            repaired_buffer = io.BytesIO()
            with open(temp_output_path, 'rb') as f:
                repaired_buffer.write(f.read())
            repaired_buffer.seek(0)
            
            print("✅ Advanced PowerPoint repair completed successfully")
            return repaired_buffer
            
        finally:
            # Clean up temp files
            try:
                os.unlink(temp_input_path)
                os.unlink(temp_output_path)
            except:
                pass
                
    except Exception as e:
        print(f"❌ Advanced PowerPoint repair failed: {e}")
        # Return original buffer if repair fails
        pptx_buffer.seek(0)
        return pptx_buffer

def repair_pptx_deep_clean(pptx_buffer):
    """
    Deep repair method - rebuilds presentation from scratch
    Most thorough but may lose some formatting
    """
    try:
        print("🔧 Starting deep PowerPoint repair...")
        
        # Reset buffer position
        pptx_buffer.seek(0)
        
        # Load the corrupted presentation
        source_prs = Presentation(pptx_buffer)
        print(f"✅ Loaded source presentation with {len(source_prs.slides)} slides")
        
        # Create a completely new presentation
        new_prs = Presentation()
        
        # Remove the default slide
        if len(new_prs.slides) > 0:
            slide_to_remove = new_prs.slides[0]
            slide_id = slide_to_remove.slide_id
            new_prs.part.drop_rel(slide_id)
            del new_prs.slides._sldIdLst[0]
        
        # Copy slides one by one
        for i, source_slide in enumerate(source_prs.slides):
            print(f"🔄 Processing slide {i + 1}...")
            
            # Add new slide with blank layout
            slide_layout = new_prs.slide_layouts[6]  # Blank layout
            new_slide = new_prs.slides.add_slide(slide_layout)
            
            # Copy content safely
            copy_slide_content_safe(source_slide, new_slide)
        
        # Save the rebuilt presentation
        repaired_buffer = io.BytesIO()
        new_prs.save(repaired_buffer)
        repaired_buffer.seek(0)
        
        print("✅ Deep PowerPoint repair completed successfully")
        return repaired_buffer
        
    except Exception as e:
        print(f"❌ Deep PowerPoint repair failed: {e}")
        # Return original buffer if repair fails
        pptx_buffer.seek(0)
        return pptx_buffer

def copy_slide_content_safe(source_slide, target_slide):
    """
    Safely copy slide content for deep repair
    """
    try:
        # Copy text content from shapes
        for i, source_shape in enumerate(source_slide.shapes):
            if hasattr(source_shape, 'text_frame') and source_shape.text_frame:
                # Add text box to target slide
                text_box = target_slide.shapes.add_textbox(
                    source_shape.left, source_shape.top, 
                    source_shape.width, source_shape.height
                )
                text_box.text = source_shape.text
                
                # Copy basic formatting
                for j, source_para in enumerate(source_shape.text_frame.paragraphs):
                    if j < len(text_box.text_frame.paragraphs):
                        text_box.text_frame.paragraphs[j].alignment = source_para.alignment
            
            elif hasattr(source_shape, 'table') and source_shape.table:
                # Copy table structure and content
                table = target_slide.shapes.add_table(
                    len(source_shape.table.rows), 
                    len(source_shape.table.columns),
                    source_shape.left, source_shape.top,
                    source_shape.width, source_shape.height
                ).table
                
                # Copy cell content
                for row_idx, source_row in enumerate(source_shape.table.rows):
                    for col_idx, source_cell in enumerate(source_row.cells):
                        if row_idx < len(table.rows) and col_idx < len(table.rows[row_idx].cells):
                            table.rows[row_idx].cells[col_idx].text = source_cell.text
        
    except Exception as e:
        print(f"⚠️ Error copying slide content: {e}")

def auto_repair_pptx(pptx_buffer, method="standard"):
    """
    Automatically repair PowerPoint file using specified method
    
    Args:
        pptx_buffer: BytesIO buffer containing the PPTX file
        method: "standard", "temp_file", or "deep_clean"
    
    Returns:
        Repaired BytesIO buffer
    """
    
    if method == "temp_file":
        return repair_pptx_with_temp_file(pptx_buffer)
    elif method == "deep_clean":
        return repair_pptx_deep_clean(pptx_buffer)
    else:
        return repair_pptx_buffer(pptx_buffer)
